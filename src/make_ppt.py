"""Generate a portfolio PowerPoint walking through the project end-to-end."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"
OUT = ROOT / "Chest_Xray_VLM_Project.pptx"

NAVY = RGBColor(0x1A, 0x3C, 0x6E)
BLUE = RGBColor(0x1A, 0x5F, 0xB4)
GRAY = RGBColor(0x44, 0x44, 0x44)
GREEN = RGBColor(0x2A, 0x9D, 0x4A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(15); r2.font.color.rgb = GRAY
    # accent line
    ln = slide.shapes.add_shape(1, Inches(0.62), Inches(1.45), Inches(2.2), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()


def bullets(slide, items, top=1.8, left=0.7, width=12.0, size=17):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = GRAY if lvl else RGBColor(0x22, 0x22, 0x22)
        if lvl == 0:
            r.font.bold = True
        p.space_after = Pt(7)


def image_slide(t, img, caption=None, sub=None):
    s = add_slide(); title(s, t, sub)
    pic = s.shapes.add_picture(str(img), Inches(0.9), Inches(1.75), width=Inches(11.5))
    if caption:
        cb = s.shapes.add_textbox(Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.6))
        p = cb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption; r.font.size = Pt(13); r.font.italic = True
        r.font.color.rgb = GRAY
    return s


# ---- Slide 1: Title ----
s = add_slide()
box = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(2.6))
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Chest X-ray Radiology Report Generation"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = NAVY
p2 = tf.add_paragraph(); r2 = p2.add_run()
r2.text = "A Vision-Language Model (VLM) built end-to-end on a single RTX 3070 (8 GB)"
r2.font.size = Pt(20); r2.font.color.rgb = BLUE
p3 = tf.add_paragraph(); p3.space_before = Pt(18); r3 = p3.add_run()
r3.text = "Oybek Valiyev"
r3.font.size = Pt(18); r3.font.bold = True
p4 = tf.add_paragraph(); r4 = p4.add_run()
r4.text = ("GitHub: github.com/varzik07/chest-xray-vlm    |    "
           "Live demo: huggingface.co/spaces/Valiye/chest-xray-report")
r4.font.size = Pt(13); r4.font.color.rgb = GRAY

# ---- Slide 2: Problem & Goal ----
s = add_slide(); title(s, "Problem & Goal")
bullets(s, [
    ("Task: given a chest X-ray, automatically generate the radiology findings report.", 0),
    ("A Vision-Language Model must SEE a medical image AND WRITE clinical text.", 1),
    ("Goals of the project:", 0),
    ("Learn how VLMs work internally (build one from scratch).", 1),
    ("Ship a real, deployed, clickable demo for a portfolio.", 1),
    ("Hard constraint: a single laptop GPU — RTX 3070, only 8 GB VRAM.", 0),
    ("Drives the choices: QLoRA, 4-bit quantization, small models, frozen encoders.", 1),
])

# ---- Slide 3: Dataset & Preprocessing ----
s = add_slide(); title(s, "Dataset & Preprocessing")
bullets(s, [
    ("Indiana University Chest X-ray collection: 3,955 reports + ~7,470 images.", 0),
    ("prepare_data.py: parse each XML report, extract FINDINGS / IMPRESSION.", 0),
    ("Clean the 'XXXX' anonymization tokens; pair each report with its image(s).", 1),
    ("Report-level train/val/test split (NOT image-level) to prevent leakage.", 0),
    ("Two views of the same patient never cross splits.", 1),
    ("Result: 5,958 train  /  734 val  /  738 test  (image, report) pairs.", 0),
])

# ---- Slide 4: Architecture ----
image_slide("Model Architecture", ASSETS / "architecture.png",
            caption="Vision encoder (frozen) -> connector (trained) -> language model writes the report.")

# ---- Slide 5: Approach ----
s = add_slide(); title(s, "Approach — Two Phases, Three Models")
bullets(s, [
    ("Phase 1 — Baseline (learn the internals):", 0),
    ("From scratch: DenseNet121 (frozen) + projector + GPT-2 decoder.", 1),
    ("Phase 2 — Ship a strong model:", 0),
    ("QLoRA fine-tune of SmolVLM-500M in 4-bit — fits in 8 GB.", 1),
    ("v3 — Headline upgrade:", 0),
    ("Swap in a CheXNet (medical) DenseNet encoder pretrained on chest X-rays.", 1),
    ("Key techniques: QLoRA, 4-bit (NF4) quantization, LoRA adapters, mixed precision, gradient checkpointing.", 0),
])

# ---- Slide 6: Results ----
s = add_slide(); title(s, "Results (test set, 738 reports)")
rows = [
    ("Model", "Encoder", "BLEU-1", "ROUGE-L", "METEOR", "Unique"),
    ("Baseline (DenseNet+GPT-2)", "ImageNet", "0.279", "0.286", "0.233", "10"),
    ("SmolVLM QLoRA v1", "SigLIP", "0.227", "0.270", "0.210", "25"),
    ("SmolVLM QLoRA v2", "SigLIP", "0.154", "0.281", "0.177", "46"),
    ("CheXNet + GPT-2 (v3)", "CheXNet (medical)", "0.343", "0.295", "0.280", "9"),
]
tbl = s.shapes.add_table(len(rows), len(rows[0]), Inches(0.7), Inches(1.9),
                         Inches(11.9), Inches(3.2)).table
for c in range(len(rows[0])):
    for rr in range(len(rows)):
        cell = tbl.cell(rr, c); cell.text = rows[rr][c]
        para = cell.text_frame.paragraphs[0]; para.font.size = Pt(14)
        if rr == 0:
            para.font.bold = True; para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        elif rr == 4:
            para.font.bold = True; cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE9, 0xF6, 0xE9)
bx = s.shapes.add_textbox(Inches(0.7), Inches(5.4), Inches(12), Inches(1.6))
p = bx.text_frame.paragraphs[0]; r = p.add_run()
r.text = "The medical CheXNet encoder won: BLEU-1 0.343 (+23% over baseline)."
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GREEN

# ---- Slide 7: Key finding ----
s = add_slide(); title(s, "Key Finding — Mode Collapse & the Metric Trap")
bullets(s, [
    ("The first QLoRA model 'looked fine' on metrics — but diversity exposed a problem:", 0),
    ("Only 25 unique reports across 738 images; it IGNORED the X-ray (language prior).", 1),
    ("Diagnosed it, traced the cause: frozen connector + class imbalance + under-training.", 0),
    ("Fix: LoRA on the vision connector + oversample abnormal cases + validation split -> diversity +84%.", 1),
    ("Deeper insight (the real headline):", 0),
    ("On this data, high BLEU comes from emitting the 'normal' template — metrics reward mode collapse.", 1),
    ("=> Real radiology AI needs clinical-efficacy metrics (CheXbert-F1), not just BLEU/ROUGE.", 1),
])

# ---- Slide 8: Quantization ----
s = add_slide(); title(s, "Quantization — 4-bit vs fp16")
bullets(s, [
    ("Benchmarked SmolVLM-500M inference (20 images):", 0),
    ("fp16: 0.97 GB VRAM, 3.6 s/image.", 1),
    ("4-bit (NF4): 0.39 GB VRAM (2.5x smaller), 5.0 s/image.", 1),
    ("Honest takeaway: 4-bit trades speed for memory.", 0),
    ("Its real value here was enabling QLoRA TRAINING to fit in 8 GB.", 1),
    ("For serving a small model that already fits, fp16 is faster.", 1),
])

# ---- Slide 9: Deployment ----
s = add_slide(); title(s, "Deployment")
bullets(s, [
    ("GitHub — full source, README, and DEPLOYMENT.md (reproducible).", 0),
    ("github.com/varzik07/chest-xray-vlm", 1),
    ("Hugging Face Spaces — live Gradio demo (free CPU, fp16).", 0),
    ("Bundled the 39 MB LoRA adapter; base model downloads at runtime.", 1),
    ("Pushed via huggingface_hub.upload_folder(); builds automatically.", 1),
    ("huggingface.co/spaces/Valiye/chest-xray-report", 1),
])

# ---- Slide 10: Demo result ----
image_slide("Live Demo — Example Result", ASSETS / "demo_result.png",
            caption="Generated report vs the radiologist's ground truth (ROUGE-L 0.90 on this case).")

# ---- Slide 11: Conclusion ----
s = add_slide(); title(s, "Summary & Skills Demonstrated")
bullets(s, [
    ("Built and compared 3 vision-language models for radiology report generation on an 8 GB GPU.", 0),
    ("Diagnosed mode collapse; improved a domain-specific encoder to +23% BLEU-1.", 0),
    ("Shipped a deployed, quantized demo with reproducible training/eval code.", 0),
    ("Skills: PyTorch, Hugging Face, QLoRA, 4-bit quantization, VLM architecture,", 0),
    ("evaluation (BLEU/ROUGE/METEOR + diversity), debugging, and deployment.", 1),
    ("Code: github.com/varzik07/chest-xray-vlm   |   Demo: huggingface.co/spaces/Valiye/chest-xray-report", 0),
])

prs.save(OUT)
print("Saved", OUT)
